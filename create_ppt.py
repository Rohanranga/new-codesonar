from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()

# Slide 1: Title Slide based on the image format
blank_slide_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_slide_layout)

# College Name
tx_box1 = slide1.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
tf1 = tx_box1.text_frame
tf1.word_wrap = True
p1 = tf1.add_paragraph()
p1.text = "Nalla Malla Reddy Engineering College"
p1.font.bold = True
p1.font.size = Pt(36)
p1.font.color.rgb = RGBColor(192, 0, 0) # Red color similar to image
p1.alignment = PP_ALIGN.CENTER

# College Subtitle
tx_box2 = slide1.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.5))
tf2 = tx_box2.text_frame
p2 = tf2.add_paragraph()
p2.text = "Autonomous Institution"
p2.font.bold = True
p2.font.size = Pt(24)
p2.alignment = PP_ALIGN.CENTER

# Project Title (center of slide)
tx_box_title = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
tf_title = tx_box_title.text_frame
tf_title.word_wrap = True
p_title = tf_title.add_paragraph()
p_title.text = "CodeSonar\nNext-Generation AI-Powered Repository Intelligence & Visualization"
p_title.font.bold = True
p_title.font.size = Pt(28)
p_title.alignment = PP_ALIGN.CENTER

# Info Box (Bottom Right)
left = Inches(5.5)
top = Inches(5)
width = Inches(4)
height = Inches(2)
tx_box3 = slide1.shapes.add_textbox(left, top, width, height)
# Add border to text box
line = tx_box3.line
line.color.rgb = RGBColor(0, 112, 192) # Blue border like image
tf3 = tx_box3.text_frame
tf3.word_wrap = True
p3 = tf3.add_paragraph()
p3.text = "By Students Names:\n1.\n2.\n3.\n4.\n\nUnder the Guidance of: Mr. Faculty Name"
p3.font.size = Pt(14)
p3.alignment = PP_ALIGN.LEFT

# Slide 2: Abstract
title_bullet_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(title_bullet_slide_layout)
slide2.shapes.title.text = "Abstract"
tf2 = slide2.placeholders[1].text_frame
tf2.text = "CodeSonar is an advanced, automated code analysis platform that transforms raw, unstructured GitHub repositories into deeply understandable, documented, and visualized architectures."
p2_1 = tf2.add_paragraph()
p2_1.text = "By combining cutting-edge Generative AI (Google Gemini) with semantic web technologies and real-time AST parsing, CodeSonar reduces developer onboarding time from weeks to mere minutes."
p2_1.level = 0

# Slide 3: Core Features
slide3 = prs.slides.add_slide(title_bullet_slide_layout)
slide3.shapes.title.text = "Core Features"
tf3 = slide3.placeholders[1].text_frame
tf3.text = "Automated Ingestion: Seamless integration with GitHub via Octokit to securely clone and traverse complex codebases."
p3_1 = tf3.add_paragraph()
p3_1.text = "Deep Contextual AI: Leverages Google Generative AI to understand intent, map dependencies, and uncover hidden architectural patterns."
p3_2 = tf3.add_paragraph()
p3_2.text = "Dynamic Rendering: Translates AI outputs into real-time interactive diagrams (Mermaid.js) and 3D Spline experiences."
p3_3 = tf3.add_paragraph()
p3_3.text = "Markdown Generation: Auto-generates exhaustive, developer-ready documentation dynamically based on real code."

# Slide 4: System Architecture Data-Flow
slide4 = prs.slides.add_slide(title_bullet_slide_layout)
slide4.shapes.title.text = "System Architecture Data-Flow"
tf4 = slide4.placeholders[1].text_frame
tf4.text = "User Provides Repo URL → Next.js Client UI"
p4_1 = tf4.add_paragraph()
p4_1.text = "Next.js Client UI → Requests Analysis from Next.js API Routes"
p4_2 = tf4.add_paragraph()
p4_2.text = "API Routes → Fetches Source Code via Octokit API connecting to GitHub"
p4_3 = tf4.add_paragraph()
p4_3.text = "API Routes → Sends Code Chunks + Context to GenKit / Google Generative AI"
p4_4 = tf4.add_paragraph()
p4_4.text = "API Routes → Aggregates & Formats via Markdown/AST Parser"
p4_5 = tf4.add_paragraph()
p4_5.text = "Next.js UI renders results using React Markdown and Mermaid.js"

# Slide 5: Value Proposition
slide5 = prs.slides.add_slide(title_bullet_slide_layout)
slide5.shapes.title.text = "Value Proposition"
tf5 = slide5.placeholders[1].text_frame
tf5.text = "Accelerated Onboarding: New engineers grasp repository layout and flow 10x faster using visual mappings."
p5_1 = tf5.add_paragraph()
p5_1.text = "Living Documentation: Never suffer from outdated READMEs. Documentation is generated directly from source on-demand."
p5_2 = tf5.add_paragraph()
p5_2.text = "Code Modernization: Identifies technical debt and anti-patterns immediately by exposing the system structure visually."

# Slide 6: Technology Stack
slide6 = prs.slides.add_slide(title_bullet_slide_layout)
slide6.shapes.title.text = "Technology Stack"
tf6 = slide6.placeholders[1].text_frame
tf6.text = "Next.js 15 (App Router) & React 19"
p6_1 = tf6.add_paragraph()
p6_1.text = "TypeScript"
p6_2 = tf6.add_paragraph()
p6_2.text = "Google Generative AI & GenKit Core"
p6_3 = tf6.add_paragraph()
p6_3.text = "Octokit (GitHub API)"
p6_4 = tf6.add_paragraph()
p6_4.text = "Tailwind CSS v4 & Framer Motion"
p6_5 = tf6.add_paragraph()
p6_5.text = "Mermaid.js 11 & React Spline"

output_path = 'r:/dup/codesonar/CodeSonar_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved successfully at {output_path}")
